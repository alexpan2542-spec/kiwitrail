"""Generate the Kiwi Trail project presentation as a .pptx file.

Designed for a 7–8 minute group project presentation, following the rubric:
content (clarity, organisation, depth), delivery, time management, and Q&A.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


# Paths to real screenshots (relative to repository root)
ASSETS_DIR = os.path.join(os.path.dirname(__file__), "..", "assets")
SHOT_HOME = os.path.join(ASSETS_DIR, "homepage.png")
SHOT_REGION = os.path.join(ASSETS_DIR, "search-waikato.png")
SHOT_DETAIL = os.path.join(ASSETS_DIR, "search.png")
SHOT_ELEVATION = os.path.join(ASSETS_DIR, "route-elevation.png")


# Kiwi-inspired colour palette
NZ_GREEN = RGBColor(0x1B, 0x5E, 0x20)        # deep forest green
NZ_GREEN_LIGHT = RGBColor(0x4C, 0xAF, 0x50)  # accent green
NZ_BROWN = RGBColor(0x5D, 0x40, 0x37)        # earthy brown
NZ_SKY = RGBColor(0x29, 0x79, 0xB5)          # sky / lake blue
NZ_GOLD = RGBColor(0xE6, 0xA8, 0x17)         # warm gold accent
TEXT_DARK = RGBColor(0x21, 0x21, 0x21)
TEXT_MUTED = RGBColor(0x55, 0x55, 0x55)
TEXT_LIGHT = RGBColor(0xFA, 0xFA, 0xFA)
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xF2)
BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)
GRID_LINE = RGBColor(0xDD, 0xDD, 0xDD)


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def add_background(slide, color):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_side_band(slide, color, width=Inches(0.3)):
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, width, prs.slide_height
    )
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()
    return band


def add_page_number(slide, number, total):
    box = slide.shapes.add_textbox(
        Inches(12.1), Inches(7.05), Inches(1.1), Inches(0.35)
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{number} / {total}"
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT_MUTED
    run.font.italic = True


def add_footer(slide, text="Kiwi Trail  |  570 Programming Project"):
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.05), Inches(10.0), Inches(0.35)
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = NZ_GREEN
    run.font.italic = True


def add_title(slide, text, subtitle=None, top=Inches(0.4)):
    box = slide.shapes.add_textbox(
        Inches(0.7), top, Inches(12.0), Inches(0.9)
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = NZ_GREEN

    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.72), top + Inches(0.78),
        Inches(1.0), Inches(0.06),
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = NZ_GREEN_LIGHT
    underline.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(2.0), top + Inches(0.22),
            Inches(10.5), Inches(0.5),
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "  •  " + subtitle
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.color.rgb = TEXT_MUTED


def add_presenter_tag(slide, name):
    """Small label in the top-right indicating which group member presents."""
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.6), Inches(0.45),
        Inches(2.55), Inches(0.42),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = NZ_GREEN
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"Presenter: {name}"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = TEXT_LIGHT


def add_bullets(slide, items, left=Inches(0.9), top=Inches(1.55),
                width=Inches(11.5), height=Inches(5.2),
                font_size=20, color=TEXT_DARK, bullet_color=NZ_GREEN_LIGHT,
                spacing=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)

    for i, item in enumerate(items):
        # Items can be (main, sub) tuples to add a smaller second line
        if isinstance(item, tuple):
            main, sub = item
        else:
            main, sub = item, None

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)

        bullet_run = p.add_run()
        bullet_run.text = "●  "
        bullet_run.font.size = Pt(font_size)
        bullet_run.font.color.rgb = bullet_color
        bullet_run.font.bold = True

        text_run = p.add_run()
        text_run.text = main
        text_run.font.size = Pt(font_size)
        text_run.font.color.rgb = color
        text_run.font.bold = True

        if sub:
            sub_p = tf.add_paragraph()
            sub_p.alignment = PP_ALIGN.LEFT
            sub_p.space_after = Pt(spacing)
            sub_p.level = 1
            sub_run = sub_p.add_run()
            sub_run.text = "      " + sub
            sub_run.font.size = Pt(font_size - 4)
            sub_run.font.color.rgb = TEXT_MUTED
    return box


def add_speaker_note(slide, note):
    slide.notes_slide.notes_text_frame.text = note


def add_card(slide, left, top, width, height, title, items,
             accent_color=NZ_GREEN, title_size=18, body_size=14):
    """A clean rounded card with coloured top stripe + title + bullet list."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = accent_color
    card.line.width = Pt(1.25)
    card.shadow.inherit = False

    # Top stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.5)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent_color
    stripe.line.fill.background()

    head_tf = stripe.text_frame
    head_tf.margin_left = Inches(0.2)
    head_tf.margin_top = Inches(0.07)
    p = head_tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(title_size)
    run.font.color.rgb = TEXT_LIGHT

    body = slide.shapes.add_textbox(
        left + Inches(0.2),
        top + Inches(0.65),
        width - Inches(0.4),
        height - Inches(0.75),
    )
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        bullet = p.add_run()
        bullet.text = "•  "
        bullet.font.size = Pt(body_size)
        bullet.font.color.rgb = accent_color
        bullet.font.bold = True
        run = p.add_run()
        run.text = item
        run.font.size = Pt(body_size)
        run.font.color.rgb = TEXT_DARK
    return card


def add_picture_with_caption(slide, image_path, left, top, width, height,
                              caption=None, fig_label=None,
                              border_color=NZ_GREEN):
    """Insert an image inside a rounded card, with optional caption below."""
    # Card background behind the image
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid(); card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = border_color; card.line.width = Pt(1.25)

    img_area_h = height
    cap_h = Inches(0)
    if caption:
        cap_h = Inches(0.85)
        img_area_h = height - cap_h - Inches(0.1)

    # Insert image, sized to fit inside the card with small padding
    pad = Inches(0.1)
    pic_left = left + pad
    pic_top = top + pad
    pic_w = width - 2 * pad
    pic_h = img_area_h - 2 * pad
    slide.shapes.add_picture(image_path, pic_left, pic_top, pic_w, pic_h)

    if caption:
        cap_box = slide.shapes.add_textbox(
            left + Inches(0.15),
            top + img_area_h,
            width - Inches(0.3),
            cap_h,
        )
        tf = cap_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.02)

        if fig_label:
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = fig_label + " "
            run.font.bold = True
            run.font.size = Pt(11)
            run.font.color.rgb = NZ_GREEN
            run2 = p.add_run()
            run2.text = caption
            run2.font.size = Pt(11)
            run2.font.color.rgb = TEXT_MUTED
        else:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = caption
            run.font.size = Pt(11)
            run.font.italic = True
            run.font.color.rgb = TEXT_MUTED


def base_slide(page_num, total_pages):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide, BG_LIGHT)
    add_side_band(slide, NZ_GREEN)
    add_footer(slide)
    add_page_number(slide, page_num, total_pages)
    return slide


# ---------------------------------------------------------------------------
# Build the presentation
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


TOTAL_SLIDES = 12   # title, problem, objective, solution, architecture,
                     # features, implementation, demo (home), demo (results),
                     # challenges, future + conclusion, Q&A


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NZ_GREEN)

# Decorative circles
c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2),
                             Inches(5), Inches(5))
c1.fill.solid(); c1.fill.fore_color.rgb = NZ_GREEN_LIGHT
c1.line.fill.background()

c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4.5),
                             Inches(6), Inches(6))
c2.fill.solid(); c2.fill.fore_color.rgb = NZ_BROWN
c2.line.fill.background()

c3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-1),
                             Inches(2.5), Inches(2.5))
c3.fill.solid(); c3.fill.fore_color.rgb = NZ_GOLD
c3.line.fill.background()

# Course tag
tag = slide.shapes.add_textbox(Inches(0.85), Inches(1.3), Inches(11), Inches(0.4))
tf = tag.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "570 Programming Project"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = NZ_GOLD

# Main title
box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.8))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Kiwi Trail"
run.font.size = Pt(90)
run.font.bold = True
run.font.color.rgb = TEXT_LIGHT

# Subtitle
box = slide.shapes.add_textbox(Inches(0.85), Inches(3.5), Inches(11.5), Inches(1.0))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "A web application for planning outdoor activities in New Zealand"
run.font.size = Pt(22)
run.font.color.rgb = TEXT_LIGHT
run.font.italic = True

# Accent line
accent = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(4.4), Inches(2.5), Inches(0.08)
)
accent.fill.solid(); accent.fill.fore_color.rgb = NZ_GOLD
accent.line.fill.background()

# Team / date placeholder
box = slide.shapes.add_textbox(Inches(0.85), Inches(4.7), Inches(11.5), Inches(2.0))
tf = box.text_frame
tf.word_wrap = True
labels = [
    "Team: [Team name]",
    "Members: [Member 1]  •  [Member 2]  •  [Member 3]  •  [Member 4]",
    "Date: [Presentation date]",
]
for i, line in enumerate(labels):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(6)
    run = p.add_run()
    run.text = line
    run.font.size = Pt(18)
    run.font.color.rgb = TEXT_LIGHT

add_speaker_note(
    slide,
    "[~30 sec] Greet the audience and introduce the team. "
    "Say: 'Hi everyone, we are Team [name] and today we will introduce "
    "Kiwi Trail — a web application that helps people plan outdoor "
    "activities in New Zealand.' Briefly mention each presenter and the "
    "section they will cover."
)


# ---------------------------------------------------------------------------
# Slide 2 — Problem & Motivation
# ---------------------------------------------------------------------------
slide = base_slide(2, TOTAL_SLIDES)
add_title(slide, "Problem & Motivation", subtitle="Why Kiwi Trail?")
add_presenter_tag(slide, "[Member 1]")

# Left: where data lives now
add_card(
    slide,
    Inches(0.9), Inches(1.7),
    Inches(5.9), Inches(2.4),
    "NZ outdoor data lives in many places",
    [
        "DOC: tracks, huts, campsites",
        "LINZ: maps and place names",
        "NIWA: weather forecasts",
        "Council sites: regional parks",
    ],
    accent_color=NZ_SKY,
    title_size=16,
    body_size=14,
)

# Right: what users need
add_card(
    slide,
    Inches(7.0), Inches(1.7),
    Inches(5.4), Inches(2.4),
    "What users actually want",
    [
        "Find places nearby quickly",
        "Filter by region & difficulty",
        "See details, map, and weather",
        "Plan trips in one place",
    ],
    accent_color=NZ_GREEN_LIGHT,
    title_size=16,
    body_size=14,
)

# Problem callout
card = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.9), Inches(4.4), Inches(11.5), Inches(2.3),
)
card.fill.solid(); card.fill.fore_color.rgb = NZ_GREEN
card.line.fill.background()
tf = card.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "The Problem"
run.font.bold = True; run.font.size = Pt(20)
run.font.color.rgb = NZ_GOLD
p = tf.add_paragraph(); p.space_before = Pt(8)
run = p.add_run()
run.text = (
    "Outdoor activity information is scattered across many sources. "
    "Users have to switch between websites to find tracks, check maps, "
    "and view the weather — making trip planning slow and frustrating."
)
run.font.size = Pt(18); run.font.color.rgb = TEXT_LIGHT

add_speaker_note(
    slide,
    "[~45 sec] Explain that New Zealand has rich outdoor data, but it is "
    "spread across DOC, LINZ, NIWA, and council websites. Give an example: "
    "'Imagine you want to go hiking this weekend — you might need 4 or 5 "
    "tabs to figure out which track, where it is, how hard it is, and what "
    "the weather will be.' This motivates building Kiwi Trail."
)


# ---------------------------------------------------------------------------
# Slide 3 — Objective & Target Users
# ---------------------------------------------------------------------------
slide = base_slide(3, TOTAL_SLIDES)
add_title(slide, "Objective & Target Users",
          subtitle="One platform for outdoor planning")
add_presenter_tag(slide, "[Member 1]")

# Objective bullets
add_bullets(
    slide,
    [
        ("Combine NZ outdoor datasets in one app",
         "Tracks, huts, campsites, weather, maps"),
        ("Help users plan trips faster",
         "Search, filter, view details, navigate"),
        ("Provide a clean and modern map-based UI",
         "Built for both desktop and mobile browsers"),
    ],
    left=Inches(0.9), top=Inches(1.6),
    width=Inches(7.5), height=Inches(5.0),
    font_size=18, spacing=4,
)

# Target users cards on the right
add_card(
    slide,
    Inches(8.7), Inches(1.6),
    Inches(3.9), Inches(1.6),
    "Local Kiwis",
    [
        "Plan weekend hikes",
        "Discover new tracks",
    ],
    accent_color=NZ_GREEN_LIGHT,
    title_size=15, body_size=13,
)
add_card(
    slide,
    Inches(8.7), Inches(3.3),
    Inches(3.9), Inches(1.6),
    "Visitors / Tourists",
    [
        "Explore NZ outdoors",
        "Quick & easy planning",
    ],
    accent_color=NZ_SKY,
    title_size=15, body_size=13,
)
add_card(
    slide,
    Inches(8.7), Inches(5.0),
    Inches(3.9), Inches(1.6),
    "Casual Campers",
    [
        "Find nearby campsites",
        "Check weather first",
    ],
    accent_color=NZ_BROWN,
    title_size=15, body_size=13,
)

add_speaker_note(
    slide,
    "[~45 sec] State the goal in one sentence: 'Make outdoor planning easier "
    "by combining NZ datasets into one web app.' Highlight three user "
    "groups — locals, visitors, and casual campers — and how each benefits."
)


# ---------------------------------------------------------------------------
# Slide 4 — Proposed Solution
# ---------------------------------------------------------------------------
slide = base_slide(4, TOTAL_SLIDES)
add_title(slide, "Proposed Solution",
          subtitle="A map-based web app with planning tools")
add_presenter_tag(slide, "[Member 2]")

# Left: home-page screenshot as a visual hook
add_picture_with_caption(
    slide,
    SHOT_HOME,
    left=Inches(0.9), top=Inches(1.6),
    width=Inches(6.5), height=Inches(5.1),
    fig_label="Preview.",
    caption=("KiwiTrail home page — search panel on the left, "
             "interactive map of New Zealand on the right."),
    border_color=NZ_GREEN_LIGHT,
)

# Right: two stacked feature lists
list_w = Inches(5.0)
list_h = Inches(2.45)
list_left = Inches(7.6)
list_top1 = Inches(1.6)
list_top2 = Inches(4.25)

# Discover
list_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    list_left, list_top1, list_w, list_h,
)
list_box.fill.solid(); list_box.fill.fore_color.rgb = BG_CARD
list_box.line.color.rgb = NZ_GREEN_LIGHT; list_box.line.width = Pt(1.25)

head = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, list_left, list_top1, list_w, Inches(0.45),
)
head.fill.solid(); head.fill.fore_color.rgb = NZ_GREEN_LIGHT
head.line.fill.background()
tf = head.text_frame; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.05)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Discover places"
run.font.bold = True; run.font.size = Pt(15)
run.font.color.rgb = TEXT_LIGHT

add_bullets(slide,
            ["Select region & difficulty",
             "Filter track / hut / campsite",
             "Fuzzy place-name search",
             "Nearby within 15 km"],
            left=list_left + Inches(0.15), top=list_top1 + Inches(0.55),
            width=list_w - Inches(0.3), height=Inches(1.8),
            font_size=13, spacing=4)

# Plan
list_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    list_left, list_top2, list_w, list_h,
)
list_box.fill.solid(); list_box.fill.fore_color.rgb = BG_CARD
list_box.line.color.rgb = NZ_SKY; list_box.line.width = Pt(1.25)

head = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, list_left, list_top2, list_w, Inches(0.45),
)
head.fill.solid(); head.fill.fore_color.rgb = NZ_SKY
head.line.fill.background()
tf = head.text_frame; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.05)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Plan your trip"
run.font.bold = True; run.font.size = Pt(15)
run.font.color.rgb = TEXT_LIGHT

add_bullets(slide,
            ["View item details",
             "Open Google Maps directions",
             "View NIWA weather",
             "Save favourites & comments"],
            left=list_left + Inches(0.15), top=list_top2 + Inches(0.55),
            width=list_w - Inches(0.3), height=Inches(1.8),
            font_size=13, spacing=4)

add_speaker_note(
    slide,
    "[~45 sec] Walk the audience through the two halves of the app: "
    "'Discover' (how users find places) and 'Plan' (what they do once "
    "they've picked one). Emphasise that the 15 km nearby search and "
    "fuzzy search make discovery very fast."
)


# ---------------------------------------------------------------------------
# Slide 5 — System Architecture
# ---------------------------------------------------------------------------
slide = base_slide(5, TOTAL_SLIDES)
add_title(slide, "System Architecture",
          subtitle="React • FastAPI • PostgreSQL")
add_presenter_tag(slide, "[Member 2]")

# Left: stack & data sources
add_card(
    slide,
    Inches(0.9), Inches(1.6),
    Inches(5.6), Inches(2.4),
    "Technology Stack",
    [
        "Frontend: React",
        "Backend: FastAPI (Python)",
        "Database: PostgreSQL",
    ],
    accent_color=NZ_GREEN,
    title_size=18, body_size=15,
)

add_card(
    slide,
    Inches(0.9), Inches(4.2),
    Inches(5.6), Inches(2.5),
    "Data Sources",
    [
        "DOC: tracks, huts, campsites",
        "LINZ: map tiles + gazetteer",
        "Region boundary data",
        "NIWA weather widget",
    ],
    accent_color=NZ_BROWN,
    title_size=18, body_size=15,
)

# Right: layered architecture diagram
diagram_left = Inches(7.0)
diagram_top = Inches(1.55)
layer_w = Inches(5.4)
layer_h = Inches(0.7)
gap = Inches(0.2)

layers = [
    ("User (browser)", NZ_GOLD),
    ("React Frontend  —  UI / map / search", NZ_GREEN_LIGHT),
    ("FastAPI Backend  —  REST APIs", NZ_SKY),
    ("PostgreSQL  —  app database", NZ_GREEN),
    ("DOC / LINZ / Region / Weather sources", NZ_BROWN),
]
y = diagram_top
prev_bottom = None
for label, color in layers:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        diagram_left, y, layer_w, layer_h,
    )
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.bold = True; run.font.size = Pt(14)
    run.font.color.rgb = TEXT_LIGHT

    if prev_bottom is not None:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            diagram_left + layer_w / 2 - Inches(0.12),
            prev_bottom,
            Inches(0.24),
            gap,
        )
        arrow.fill.solid(); arrow.fill.fore_color.rgb = NZ_GREEN
        arrow.line.fill.background()
    prev_bottom = y + layer_h
    y = prev_bottom + gap

add_speaker_note(
    slide,
    "[~45 sec] Explain the three-tier architecture briefly. "
    "React talks to FastAPI, FastAPI talks to PostgreSQL, and we "
    "imported DOC, LINZ, regional, and NIWA data into the database. "
    "Mention why we chose this stack (fast, light, easy to deploy)."
)


# ---------------------------------------------------------------------------
# Slide 6 — Key Features
# ---------------------------------------------------------------------------
slide = base_slide(6, TOTAL_SLIDES)
add_title(slide, "Key Features", subtitle="What users can do")
add_presenter_tag(slide, "[Member 3]")

features = [
    ("Region selection", NZ_GREEN),
    ("Difficulty filter", NZ_GREEN_LIGHT),
    ("Type filter (track / hut / camp)", NZ_SKY),
    ("Fuzzy search", NZ_BROWN),
    ("Nearby within 15 km", NZ_GOLD),
    ("Interactive map", NZ_GREEN),
    ("Detail pages", NZ_GREEN_LIGHT),
    ("Google Maps directions", NZ_SKY),
    ("Weather forecast", NZ_BROWN),
    ("User comments", NZ_GOLD),
    ("Favourites list", NZ_GREEN),
    ("Mobile-friendly UI", NZ_SKY),
]

cols = 3
pill_w = Inches(3.9)
pill_h = Inches(1.05)
gap_x = Inches(0.25)
gap_y = Inches(0.22)
start_left = Inches(0.9)
start_top = Inches(1.6)

for i, (text, color) in enumerate(features):
    row, col = divmod(i, cols)
    left = start_left + col * (pill_w + gap_x)
    top = start_top + row * (pill_h + gap_y)

    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, pill_w, pill_h
    )
    pill.fill.solid(); pill.fill.fore_color.rgb = BG_CARD
    pill.line.color.rgb = color; pill.line.width = Pt(1.5)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.18), pill_h,
    )
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    tf = pill.text_frame
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.1)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = True; run.font.size = Pt(15)
    run.font.color.rgb = TEXT_DARK

add_speaker_note(
    slide,
    "[~45 sec] Walk the audience through the headline features at a high "
    "level — don't list every single one. Group them: 'discovery features' "
    "(search / filter / nearby), 'detail features' (page / weather / map "
    "link), and 'user features' (comments / favourites). Tell the audience "
    "they'll see these in action in the demo."
)


# ---------------------------------------------------------------------------
# Slide 7 — Implementation
# ---------------------------------------------------------------------------
slide = base_slide(7, TOTAL_SLIDES)
add_title(slide, "Implementation",
          subtitle="Frontend • Backend • Database")
add_presenter_tag(slide, "[Member 3]")

card_w = Inches(4.0)
card_h = Inches(4.8)
card_top = Inches(1.7)
gap = Inches(0.2)

add_card(
    slide,
    Inches(0.9), card_top, card_w, card_h,
    "Frontend (React)",
    [
        "Interactive map view",
        "Search & filter panel",
        "Item detail pages",
        "Comments section",
        "Favourites list",
    ],
    accent_color=NZ_GREEN_LIGHT,
    title_size=17, body_size=14,
)

add_card(
    slide,
    Inches(0.9) + card_w + gap, card_top, card_w, card_h,
    "Backend (FastAPI)",
    [
        "REST API endpoints",
        "Region / difficulty filters",
        "Nearby (15 km) search",
        "Fuzzy search logic",
        "Comments & favourites",
    ],
    accent_color=NZ_SKY,
    title_size=17, body_size=14,
)

add_card(
    slide,
    Inches(0.9) + 2 * (card_w + gap), card_top, card_w, card_h,
    "Database (PostgreSQL)",
    [
        "tracks, huts, campsites",
        "regions",
        "users",
        "comments",
        "favourites",
    ],
    accent_color=NZ_BROWN,
    title_size=17, body_size=14,
)

# Workflow strip at the bottom
strip_top = Inches(6.65)
steps = ["Collect data", "Clean & import", "Expose APIs", "Show on map"]
step_w = Inches(2.9)
step_gap = Inches(0.1)
start_left = Inches(0.9)
for i, label in enumerate(steps):
    left = start_left + i * (step_w + step_gap)
    box = slide.shapes.add_shape(
        MSO_SHAPE.PENTAGON, left, strip_top, step_w, Inches(0.55)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = [NZ_GREEN, NZ_GREEN_LIGHT, NZ_SKY, NZ_BROWN][i]
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.1); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{i+1}. {label}"
    run.font.bold = True; run.font.size = Pt(12)
    run.font.color.rgb = TEXT_LIGHT

add_speaker_note(
    slide,
    "[~45 sec] Explain how the project is divided across the three layers. "
    "Mention what each team member worked on if relevant. Briefly call out "
    "the data pipeline at the bottom: collect → clean → import → expose → "
    "display."
)


# ---------------------------------------------------------------------------
# Slide 8 — Demo (Home page)
# ---------------------------------------------------------------------------
slide = base_slide(8, TOTAL_SLIDES)
add_title(slide, "Demo — Home Page", subtitle="Search panel + region map")
add_presenter_tag(slide, "[Member 4]")

# Big screenshot on the left
add_picture_with_caption(
    slide,
    SHOT_REGION,
    left=Inches(0.9), top=Inches(1.6),
    width=Inches(8.0), height=Inches(5.1),
    fig_label="Figure 1.",
    caption=("The KiwiTrail home page and search panel. Users choose a region, "
             "set the track difficulty, pick the item types (tracks, huts, "
             "campsites, weather stations) and optionally use a fuzzy search."),
    border_color=NZ_GREEN,
)

# Right: explanation of what's on screen
add_card(
    slide,
    Inches(9.1), Inches(1.6),
    Inches(3.3), Inches(5.1),
    "Search panel features",
    [
        "Region selector",
        "Track difficulty filter",
        "Type filter — tracks, huts, campsites, weather stations",
        "Fuzzy place-name search",
        "Region boundary drawn on the map",
        "Live latitude / longitude readout",
    ],
    accent_color=NZ_GREEN_LIGHT,
    title_size=16, body_size=13,
)

add_speaker_note(
    slide,
    "[~45–60 sec] Open the demo with the home page. Walk through the search "
    "panel: pick a region (e.g. Waikato), choose difficulty, tick the item "
    "types you want, and optionally type a rough place name in fuzzy search. "
    "Point out that the map highlights the selected region's boundary so the "
    "user immediately sees the area they will explore."
)


# ---------------------------------------------------------------------------
# Slide 9 — Demo (Search result + 3D elevation)
# ---------------------------------------------------------------------------
slide = base_slide(9, TOTAL_SLIDES)
add_title(slide, "Demo — Search Result",
          subtitle="Track on map, detail panel, 3D elevation")
add_presenter_tag(slide, "[Member 4]")

# Top-left: track on map + detail panel screenshot
add_picture_with_caption(
    slide,
    SHOT_DETAIL,
    left=Inches(0.9), top=Inches(1.55),
    width=Inches(6.0), height=Inches(3.6),
    fig_label="Figure 2a.",
    caption=("Track shown on the 2D map with the detail panel: photo, "
             "description, difficulty, duration, type, and 'Track Details' link."),
    border_color=NZ_SKY,
)

# Top-right: elevation screenshot
add_picture_with_caption(
    slide,
    SHOT_ELEVATION,
    left=Inches(7.1), top=Inches(1.55),
    width=Inches(5.3), height=Inches(3.6),
    fig_label="Figure 2b.",
    caption=("Route view with elevation profile — distance vs. elevation "
             "for the selected track (route 9 / 12, length 9329 m)."),
    border_color=NZ_GREEN,
)

# Bottom: action panel summary
actions = [
    ("Directions", "Open in Google Maps", NZ_SKY),
    ("Weather", "NIWA forecast widget", NZ_GREEN_LIGHT),
    ("Comments", "Read & post comments", NZ_BROWN),
    ("Favourites", "Save for later", NZ_GOLD),
]
btn_w = Inches(2.85)
btn_h = Inches(1.4)
btn_gap = Inches(0.15)
btn_top = Inches(5.3)
btn_start_left = Inches(0.9)
for i, (label, desc, color) in enumerate(actions):
    left = btn_start_left + i * (btn_w + btn_gap)
    btn = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, btn_top, btn_w, btn_h
    )
    btn.fill.solid(); btn.fill.fore_color.rgb = BG_CARD
    btn.line.color.rgb = color; btn.line.width = Pt(1.5)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, btn_top, btn_w, Inches(0.4)
    )
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.bold = True; run.font.size = Pt(14)
    run.font.color.rgb = TEXT_LIGHT

    body = slide.shapes.add_textbox(
        left + Inches(0.15), btn_top + Inches(0.5),
        btn_w - Inches(0.3), btn_h - Inches(0.55),
    )
    tf = body.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = desc
    run.font.size = Pt(12); run.font.color.rgb = TEXT_DARK

# Footer caption for action row
cap = slide.shapes.add_textbox(
    Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.3)
)
tf = cap.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = ("Detail panel actions — Directions, Weather, Comments, "
            "and Favourites — for each selected track.")
run.font.size = Pt(11); run.font.italic = True
run.font.color.rgb = TEXT_MUTED

add_speaker_note(
    slide,
    "[~60–75 sec] Pick a track on the map (e.g. Karioi Track or North-South "
    "Track). Show the detail panel: photo, description, difficulty, duration, "
    "type, and the 'View Official Details' link. Then open the 3D elevation "
    "view to show distance-vs-elevation. Finish by demonstrating the four "
    "action buttons: Directions (Google Maps), Weather (NIWA), Comments, "
    "and Favourites."
)


# ---------------------------------------------------------------------------
# Slide 10 — Challenges & Solutions
# ---------------------------------------------------------------------------
slide = base_slide(10, TOTAL_SLIDES)
add_title(slide, "Challenges & Solutions",
          subtitle="What we learned along the way")
add_presenter_tag(slide, "[Member 4]")

challenges = [
    ("Data from multiple sources",
     "Cleaned and unified DOC, LINZ and region data in PostgreSQL.",
     NZ_GREEN),
    ("Showing geographic data on a map",
     "Used LINZ map tiles and lat/long coordinates for visual display.",
     NZ_SKY),
    ("Helping users find places fast",
     "Added fuzzy search, filters, and 15 km nearby search.",
     NZ_GREEN_LIGHT),
    ("Better trip-planning experience",
     "Added weather, Google Maps links, comments and favourites.",
     NZ_BROWN),
]

card_w = Inches(5.85)
card_h = Inches(2.4)
gap_x = Inches(0.25)
gap_y = Inches(0.25)
start_left = Inches(0.9)
start_top = Inches(1.6)

for i, (chal, sol, color) in enumerate(challenges):
    row, col = divmod(i, 2)
    left = start_left + col * (card_w + gap_x)
    top = start_top + row * (card_h + gap_y)

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h
    )
    card.fill.solid(); card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = color; card.line.width = Pt(1.5)

    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, card_w, Inches(0.55)
    )
    strip.fill.solid(); strip.fill.fore_color.rgb = color
    strip.line.fill.background()

    head_tf = strip.text_frame
    head_tf.margin_left = Inches(0.2); head_tf.margin_top = Inches(0.07)
    p = head_tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Challenge {i + 1}:  {chal}"
    run.font.bold = True; run.font.size = Pt(14)
    run.font.color.rgb = TEXT_LIGHT

    body = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.7),
        card_w - Inches(0.4), card_h - Inches(0.8),
    )
    tf = body.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Solution:  "
    run.font.bold = True; run.font.size = Pt(14)
    run.font.color.rgb = color
    run2 = p.add_run()
    run2.text = sol
    run2.font.size = Pt(14); run2.font.color.rgb = TEXT_DARK

add_speaker_note(
    slide,
    "[~45 sec] Don't read every card. Pick the most interesting one or two "
    "and tell a short story — e.g. 'integrating data from DOC and LINZ was "
    "tricky because formats and coordinate systems differed; we wrote import "
    "scripts that normalise everything before inserting into PostgreSQL.'"
)


# ---------------------------------------------------------------------------
# Slide 11 — Future Work & Conclusion
# ---------------------------------------------------------------------------
slide = base_slide(11, TOTAL_SLIDES)
add_title(slide, "Future Work & Conclusion",
          subtitle="Where Kiwi Trail can go next")
add_presenter_tag(slide, "[Member 4]")

# Left: future work
add_card(
    slide,
    Inches(0.9), Inches(1.6),
    Inches(6.0), Inches(5.1),
    "Future Improvements",
    [
        "User login & profiles",
        "Trip planning itinerary",
        "Offline map support",
        "Safety alerts & track condition updates",
        "User-uploaded photos",
        "Mobile app version",
        "Personalised recommendations",
    ],
    accent_color=NZ_SKY,
    title_size=18, body_size=15,
)

# Right: conclusion / impact
concl = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(7.1), Inches(1.6),
    Inches(5.3), Inches(5.1),
)
concl.fill.solid(); concl.fill.fore_color.rgb = NZ_GREEN
concl.line.fill.background()
tf = concl.text_frame
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
tf.margin_top = Inches(0.25); tf.word_wrap = True

p = tf.paragraphs[0]
run = p.add_run()
run.text = "Conclusion"
run.font.bold = True; run.font.size = Pt(22)
run.font.color.rgb = NZ_GOLD

for line in [
    "Kiwi Trail brings NZ outdoor data together in one place.",
    "Search, filter, map, weather, navigation — all in one app.",
    "Helps both locals and visitors plan trips quickly & safely.",
]:
    p = tf.add_paragraph()
    p.space_before = Pt(10)
    run = p.add_run()
    run.text = "•  " + line
    run.font.size = Pt(16); run.font.color.rgb = TEXT_LIGHT

p = tf.add_paragraph(); p.space_before = Pt(18)
run = p.add_run()
run.text = "Thank you!"
run.font.bold = True; run.font.size = Pt(22)
run.font.color.rgb = NZ_GOLD

add_speaker_note(
    slide,
    "[~45 sec] Summarise the project in 2–3 sentences. Mention 1–2 future "
    "directions that excite you. End with a warm 'thank you' and invite "
    "questions from the audience."
)


# ---------------------------------------------------------------------------
# Slide 12 — Q&A
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NZ_GREEN)

# Decorative shapes
c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(-1.5),
                             Inches(4.5), Inches(4.5))
c1.fill.solid(); c1.fill.fore_color.rgb = NZ_GREEN_LIGHT
c1.line.fill.background()
c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(5),
                             Inches(5), Inches(5))
c2.fill.solid(); c2.fill.fore_color.rgb = NZ_BROWN
c2.line.fill.background()

# Big "?"
qmark = slide.shapes.add_textbox(Inches(0.8), Inches(1.0),
                                   Inches(12), Inches(3.5))
tf = qmark.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Q & A"
run.font.size = Pt(120); run.font.bold = True
run.font.color.rgb = TEXT_LIGHT

# Subtitle
sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.4),
                                 Inches(12), Inches(1.0))
tf = sub.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "We'd love to hear your questions and feedback."
run.font.size = Pt(22); run.font.italic = True
run.font.color.rgb = TEXT_LIGHT

# Project tag
tag = slide.shapes.add_textbox(Inches(0.8), Inches(5.4),
                                 Inches(12), Inches(0.6))
tf = tag.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Kiwi Trail  •  570 Programming Project"
run.font.size = Pt(16); run.font.bold = True
run.font.color.rgb = NZ_GOLD

add_speaker_note(
    slide,
    "[Remaining time / discussion] Open the floor for questions. Be ready "
    "to answer about: tech stack choices, data sources, how nearby search "
    "works, scalability, and future plans. If you don't know an answer, "
    "say so honestly and offer to follow up."
)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
output_path = "Kiwi_Trail_Presentation.pptx"
prs.save(output_path)
print(f"Saved presentation to {output_path}")
print(f"Total slides: {len(prs.slides)}")
