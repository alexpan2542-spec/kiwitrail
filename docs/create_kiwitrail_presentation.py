from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "kiwitrail_project_presentation.pptx"
ASSETS = ROOT / "assets"

FONT = "Microsoft YaHei"
FONT_EN = "Aptos"

GREEN = RGBColor(31, 122, 77)
DARK_GREEN = RGBColor(20, 82, 55)
BLUE = RGBColor(30, 92, 160)
SKY = RGBColor(214, 234, 248)
ORANGE = RGBColor(242, 147, 57)
TEXT = RGBColor(35, 45, 56)
MUTED = RGBColor(96, 112, 128)
LIGHT = RGBColor(246, 249, 247)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(210, 220, 216)


def set_run(run, size=20, color=TEXT, bold=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_bg(slide, color=LIGHT):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Inches(13.333),
        Inches(7.5),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def add_title(slide, title, subtitle=None, section=None):
    if section:
        box = slide.shapes.add_textbox(Inches(0.62), Inches(0.35), Inches(3.2), Inches(0.32))
        p = box.text_frame.paragraphs[0]
        p.text = section
        p.alignment = PP_ALIGN.LEFT
        set_run(p.runs[0], 10, GREEN, True, FONT_EN)

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.62), Inches(7.3), Inches(0.62))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_run(p.runs[0], 26, DARK_GREEN, True)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.62), Inches(1.22), Inches(8.2), Inches(0.42))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        set_run(p.runs[0], 13, MUTED)


def add_footer(slide, page):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6),
        Inches(7.05),
        Inches(12.1),
        Inches(0.01),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.62), Inches(7.08), Inches(4.5), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = "KiwiTrail | Full-stack Geospatial Web Application"
    set_run(p.runs[0], 8, MUTED, False, FONT_EN)

    num = slide.shapes.add_textbox(Inches(12.25), Inches(7.08), Inches(0.5), Inches(0.25))
    p = num.text_frame.paragraphs[0]
    p.text = f"{page:02d}"
    p.alignment = PP_ALIGN.RIGHT
    set_run(p.runs[0], 8, MUTED, False, FONT_EN)


def add_bullets(slide, x, y, w, h, bullets, size=14, color=TEXT, line_spacing=1.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = line_spacing
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, title, body, accent=GREEN, number=None):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(225, 233, 229)
    card.line.width = Pt(1)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.12),
        Inches(y + 0.16),
        Inches(0.12),
        Inches(h - 0.32),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    if number is not None:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + 0.36),
            Inches(y + 0.22),
            Inches(0.42),
            Inches(0.42),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()
        t = slide.shapes.add_textbox(Inches(x + 0.36), Inches(y + 0.28), Inches(0.42), Inches(0.25))
        p = t.text_frame.paragraphs[0]
        p.text = str(number)
        p.alignment = PP_ALIGN.CENTER
        set_run(p.runs[0], 10, WHITE, True, FONT_EN)
        tx = x + 0.86
    else:
        tx = x + 0.42

    title_box = slide.shapes.add_textbox(Inches(tx), Inches(y + 0.18), Inches(w - (tx - x) - 0.25), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    set_run(p.runs[0], 14, DARK_GREEN, True)

    body_box = slide.shapes.add_textbox(Inches(x + 0.42), Inches(y + 0.62), Inches(w - 0.65), Inches(h - 0.8))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = body
    p.line_spacing = 1.05
    set_run(p.runs[0], 10.5, MUTED)
    return card


def fit_picture(slide, path, x, y, max_w, max_h, rounded=False):
    img = Image.open(path)
    iw, ih = img.size
    scale = min(max_w / iw, max_h / ih)
    w = iw * scale
    h = ih * scale
    left = x + (max_w - w) / 2
    top = y + (max_h - h) / 2
    pic = slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(w), Inches(h))
    pic.line.color.rgb = RGBColor(225, 232, 230)
    pic.line.width = Pt(1)
    if rounded:
        # PowerPoint keeps the image rectangular; the surrounding white panel gives it a card feel.
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x - 0.05),
            Inches(y - 0.05),
            Inches(max_w + 0.1),
            Inches(max_h + 0.1),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = RGBColor(225, 232, 230)
        slide.shapes._spTree.remove(panel._element)
        slide.shapes._spTree.insert(pic._element.getparent().index(pic._element), panel._element)
    return pic


def add_badge(slide, x, y, text, color=GREEN):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.25), Inches(0.34))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.06), Inches(1.25), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_run(p.runs[0], 8.5, WHITE, True, FONT_EN)


def add_flow(slide, labels, x, y, w, h, colors):
    gap = 0.25
    box_w = (w - gap * (len(labels) - 1)) / len(labels)
    centers = []
    for i, label in enumerate(labels):
        bx = x + i * (box_w + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bx),
            Inches(y),
            Inches(box_w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.fill.background()
        tx = slide.shapes.add_textbox(Inches(bx + 0.08), Inches(y + 0.18), Inches(box_w - 0.16), Inches(h - 0.25))
        p = tx.text_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        set_run(p.runs[0], 12, WHITE, True)
        centers.append((bx + box_w, y + h / 2))
        if i < len(labels) - 1:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(bx + box_w + 0.03),
                Inches(y + h / 2),
                Inches(bx + box_w + gap - 0.03),
                Inches(y + h / 2),
            )
            conn.line.color.rgb = LINE
            conn.line.width = Pt(2)
            conn.line.end_arrowhead = True


def add_architecture_node(slide, x, y, w, h, title, lines, color):
    node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    node.fill.solid()
    node.fill.fore_color.rgb = WHITE
    node.line.color.rgb = color
    node.line.width = Pt(1.5)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.38))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.08), Inches(w - 0.2), Inches(0.22))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    set_run(p.runs[0], 10.5, WHITE, True)
    add_bullets(slide, x + 0.18, y + 0.55, w - 0.35, h - 0.65, lines, size=9.2, color=TEXT)


def new_slide(prs, page, title, subtitle=None, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle, section)
    add_footer(slide, page)
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(239, 247, 241))
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.25), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = DARK_GREEN
    accent.line.fill.background()
    title = slide.shapes.add_textbox(Inches(0.62), Inches(0.82), Inches(3.15), Inches(1.45))
    p = title.text_frame.paragraphs[0]
    p.text = "KiwiTrail\n项目展示"
    p.line_spacing = 0.9
    set_run(p.runs[0], 31, WHITE, True)
    sub = slide.shapes.add_textbox(Inches(0.68), Inches(2.55), Inches(2.95), Inches(0.8))
    p = sub.text_frame.paragraphs[0]
    p.text = "新西兰户外目的地发现平台\nReact · FastAPI · PostGIS"
    p.line_spacing = 1.1
    set_run(p.runs[0], 13, RGBColor(219, 239, 227))
    add_badge(slide, 0.72, 3.65, "Geospatial", ORANGE)
    add_badge(slide, 2.1, 3.65, "Full-stack", BLUE)
    fit_picture(slide, ASSETS / "homepage.png", 4.65, 0.55, 7.9, 5.8, rounded=True)
    add_card(
        slide,
        4.8,
        6.0,
        7.55,
        0.72,
        "一句话定位",
        "把 DOC 开放数据、地形海拔、天气和导航整合到交互式地图中，帮助用户更快规划徒步、住宿和露营。",
        GREEN,
    )
    add_footer(slide, 1)

    # 2. Background
    slide = new_slide(
        prs,
        2,
        "项目背景与目标",
        "围绕户外出行规划，把“找到地点、理解路线、评估环境、快速出发”串成一个闭环。",
        "01 CONTEXT",
    )
    add_card(slide, 0.7, 1.95, 3.8, 1.55, "用户痛点", "信息分散在 DOC 页面、地图、天气站和导航工具之间；路线难度、海拔和天气需要反复切换查询。", ORANGE)
    add_card(slide, 4.85, 1.95, 3.8, 1.55, "产品目标", "以地图为主界面，按区域、难度、类型和地点名称快速定位 Tracks、Huts、Campsites 与 Weather Stations。", GREEN)
    add_card(slide, 9.0, 1.95, 3.1, 1.55, "技术目标", "用 PostGIS 进行空间查询，用 FastAPI 提供稳定接口，用 React/Leaflet 交付流畅交互体验。", BLUE)
    add_flow(
        slide,
        ["探索区域", "筛选地点", "查看详情", "评估天气/海拔", "导航出发"],
        1.0,
        4.35,
        11.25,
        0.82,
        [GREEN, BLUE, ORANGE, DARK_GREEN, RGBColor(88, 119, 169)],
    )
    add_bullets(
        slide,
        1.05,
        5.62,
        10.9,
        0.8,
        [
            "面向徒步者、露营者和新西兰户外旅行规划场景；强调“地图优先、数据可信、操作直接”。",
            "项目已部署：前端 Vercel，后端 DigitalOcean Droplet；README 中提供在线 Demo 链接。",
        ],
        size=12,
    )

    # 3. Feature overview
    slide = new_slide(prs, 3, "核心功能总览", "从发现到行动：搜索、地图浏览、路线分析、天气和用户互动。", "02 FEATURES")
    cards = [
        ("空间搜索", "按 New Zealand 区域边界、地点类型、路线难度和地名模糊搜索进行筛选。", GREEN),
        ("交互式地图", "支持 Topo、Aerial、OSM、CARTO、Esri 等底图，并在地图上展示不同类型标记。", BLUE),
        ("路线详情", "展示路线 GeoJSON、长度、最高/最低海拔与 elevation profile。", ORANGE),
        ("行动入口", "一键打开 Google Maps、查看天气面板、评论、收藏和官方详情。", DARK_GREEN),
    ]
    for i, (t, b, c) in enumerate(cards):
        add_card(slide, 0.75 + (i % 2) * 6.0, 1.75 + (i // 2) * 1.65, 5.45, 1.2, t, b, c, i + 1)
    fit_picture(slide, ASSETS / "search.png", 1.0, 5.0, 5.3, 1.4, rounded=True)
    fit_picture(slide, ASSETS / "route-elevation.png", 7.0, 5.0, 4.8, 1.4, rounded=True)

    # 4. User journey
    slide = new_slide(prs, 4, "用户使用流程", "用户不需要理解后端数据结构，只需要围绕地图完成探索和决策。", "03 USER FLOW")
    add_flow(
        slide,
        ["选择区域\n或输入地名", "选择类型\n与难度", "提交搜索\n地图展示结果", "点击标记\n查看详情", "天气/评论\n收藏/导航"],
        0.75,
        1.85,
        11.85,
        1.0,
        [GREEN, BLUE, ORANGE, DARK_GREEN, RGBColor(112, 96, 171)],
    )
    fit_picture(slide, ASSETS / "search-waikato.png", 0.85, 3.35, 5.4, 2.55, rounded=True)
    add_card(slide, 6.65, 3.34, 5.45, 0.9, "区域搜索", "前端请求 /regions/{region_code} 获取 GeoJSON，并用 Leaflet 自动 fitBounds 到区域范围。", GREEN)
    add_card(slide, 6.65, 4.43, 5.45, 0.9, "地点模糊搜索", "RapidFuzz 匹配 Gazetteer 地名，PostGIS 生成 15km buffer，再统计圈内地点数量。", BLUE)
    add_card(slide, 6.65, 5.52, 5.45, 0.9, "结果探索", "用户点击地图标记后，详情面板与路线 GeoJSON 叠加层同步打开。", ORANGE)

    # 5. Frontend
    slide = new_slide(prs, 5, "前端实现亮点", "React + TypeScript + Vite 构建地图主界面，组件按搜索、详情、天气、评论、账户拆分。", "04 FRONTEND")
    fit_picture(slide, ASSETS / "homepage.png", 0.75, 1.65, 6.1, 4.25, rounded=True)
    add_card(slide, 7.15, 1.65, 4.95, 0.95, "地图体验", "React Leaflet 管理 MapContainer、Marker、GeoJSON、LayersControl 与 Tooltip。", GREEN)
    add_card(slide, 7.15, 2.8, 4.95, 0.95, "侧边面板", "搜索面板可折叠；详情、天气和评论面板互斥展示，减少地图遮挡。", BLUE)
    add_card(slide, 7.15, 3.95, 4.95, 0.95, "用户功能", "Google 登录上下文支撑收藏；评论支持评分、文本和图片上传。", ORANGE)
    add_card(slide, 7.15, 5.1, 4.95, 0.95, "可视化", "路线详情页用 Recharts 展示海拔剖面，地图上突出显示路线几何。", DARK_GREEN)

    # 6. Backend
    slide = new_slide(prs, 6, "后端 API 与服务分层", "FastAPI 暴露 REST 接口，Repository 层负责 SQL/PostGIS 查询，Service 层组织搜索逻辑。", "05 BACKEND")
    add_architecture_node(slide, 0.85, 1.75, 3.3, 1.65, "FastAPI App", ["/search/items", "/search/gazetteer", "/track-routes/{id}", "/comments", "/favourites"], GREEN)
    add_architecture_node(slide, 5.0, 1.75, 3.3, 1.65, "Services", ["search_items_service", "search_gazetteer_service", "RapidFuzz + PostGIS buffer", "startup preload"], BLUE)
    add_architecture_node(slide, 9.15, 1.75, 3.3, 1.65, "Repositories", ["track / hut / campsite", "weather station", "comment / favourite", "region"], ORANGE)
    add_flow(slide, ["HTTP Request", "Pydantic Schema", "Service", "Repository", "PostGIS SQL"], 1.0, 4.25, 11.2, 0.82, [GREEN, BLUE, DARK_GREEN, ORANGE, RGBColor(99, 126, 145)])
    add_bullets(
        slide,
        1.0,
        5.55,
        10.9,
        0.85,
        [
            "CORS 通过 ALLOWED_ORIGINS 配置；评论图片挂载到 /static，文件名使用 uuid 防冲突。",
            "接口返回统一的地图 item 结构：id、name、type、lat/lng、thumbnail、source_page_url、weather_url 等。",
        ],
        12,
    )

    # 7. Data pipeline
    slide = new_slide(prs, 7, "数据管线与空间数据模型", "把公开数据清洗成可查询的 PostGIS 表，再由 API 面向地图端输出轻量 JSON。", "06 DATA")
    sources = ["NZ DOC\nTracks", "DOC Huts\n& Campsites", "NZ Gazetteer", "NZ 8m DEM", "NIWA / MetService"]
    for i, s in enumerate(sources):
        add_card(slide, 0.7 + i * 2.45, 1.65, 2.1, 1.0, s, "公开数据源", [GREEN, BLUE, ORANGE, DARK_GREEN, RGBColor(108, 91, 123)][i])
    add_flow(slide, ["下载/清洗", "格式转换", "PostgreSQL + PostGIS", "FastAPI 查询", "React 地图"], 1.0, 3.35, 11.25, 0.8, [GREEN, BLUE, ORANGE, DARK_GREEN, RGBColor(88, 119, 169)])
    add_card(slide, 1.0, 4.85, 3.45, 1.08, "空间字段", "Tracks 使用 LINESTRING / route geometry；Huts、Weather Stations 使用 POINT；Campsites 支持区域或点位。", GREEN)
    add_card(slide, 4.95, 4.85, 3.45, 1.08, "核心查询", "ST_Intersects、ST_PointOnSurface、ST_AsGeoJSON、ST_Buffer、ST_SetSRID。", BLUE)
    add_card(slide, 8.9, 4.85, 3.1, 1.08, "ETL 脚本", "import_tracks / import_huts / import_campsites / import_8m_dem 等脚本支撑数据入库。", ORANGE)

    # 8. Elevation
    slide = new_slide(prs, 8, "路线海拔剖面", "用 NZ 8m DEM 数据对路线按距离采样，形成可视化 elevation profile。", "07 GEO ANALYTICS")
    fit_picture(slide, ASSETS / "route-elevation.png", 0.75, 1.55, 6.2, 4.25, rounded=True)
    add_card(slide, 7.25, 1.7, 4.85, 0.92, "路线长度", "将 WGS84 路线转换到 NZTM2000 计算米级长度。", GREEN)
    add_card(slide, 7.25, 2.83, 4.85, 0.92, "DEM 采样", "沿路线每 10m 插值采样，并从相交 DEM tiles 读取高程值。", BLUE)
    add_card(slide, 7.25, 3.96, 4.85, 0.92, "数据写回", "保存 length_m、elev_min、elev_max、elevation_profile 等字段。", ORANGE)
    add_card(slide, 7.25, 5.09, 4.85, 0.92, "前端呈现", "详情页将路线几何与海拔曲线结合，帮助用户评估坡度和体力要求。", DARK_GREEN)

    # 9. Search detail
    slide = new_slide(prs, 9, "搜索能力设计", "常规区域过滤与地名模糊搜索并存，覆盖“我知道区域”和“我只记得地名”两类场景。", "08 SEARCH")
    add_architecture_node(slide, 0.85, 1.65, 3.25, 1.55, "区域/类型搜索", ["Region GeoJSON", "Track difficulty", "Tracks/Huts/Campsites", "Weather Station"], GREEN)
    add_architecture_node(slide, 5.0, 1.65, 3.25, 1.55, "Fuzzy Search", ["RapidFuzz WRatio", "Gazetteer preload", "score cutoff 65", "Top N candidates"], BLUE)
    add_architecture_node(slide, 9.15, 1.65, 3.25, 1.55, "PostGIS Circle", ["ST_MakePoint", "ST_Buffer geography", "15km radius", "items_len count"], ORANGE)
    fit_picture(slide, ASSETS / "search.png", 0.95, 3.75, 5.2, 2.25, rounded=True)
    add_bullets(
        slide,
        6.75,
        4.0,
        4.8,
        1.45,
        [
            "常规搜索直接调用 /search/items，根据 selected_area 与类型开关聚合多个 repository 查询结果。",
            "模糊搜索调用 /search/gazetteer，返回候选地点、匹配分、GeoJSON 圆形范围和圈内结果数量。",
            "用户选择候选后，前端再用该 GeoJSON 圆形范围加载地图 item。",
        ],
        12,
    )

    # 10. Deployment architecture
    slide = new_slide(prs, 10, "系统架构与部署", "前端静态部署，后端容器化运行，数据库承担空间计算与持久化。", "09 ARCHITECTURE")
    add_architecture_node(slide, 0.8, 1.7, 2.35, 1.4, "Browser", ["React UI", "Leaflet map", "User actions"], GREEN)
    add_architecture_node(slide, 3.65, 1.7, 2.35, 1.4, "Vercel", ["Static assets", "VITE_BACKEND_URL", "CDN delivery"], BLUE)
    add_architecture_node(slide, 6.5, 1.7, 2.35, 1.4, "DigitalOcean", ["Docker", "FastAPI", "Uvicorn"], ORANGE)
    add_architecture_node(slide, 9.35, 1.7, 2.35, 1.4, "PostgreSQL", ["PostGIS", "GeoJSON", "Spatial index"], DARK_GREEN)
    add_flow(slide, ["Client", "Frontend", "REST API", "Repository SQL", "PostGIS"], 1.0, 3.65, 10.9, 0.75, [GREEN, BLUE, ORANGE, DARK_GREEN, RGBColor(99, 126, 145)])
    add_card(slide, 1.1, 5.05, 3.2, 0.95, "外部地图/导航", "LINZ basemaps、OSM、CARTO、Esri 作为底图；Google Maps 用于导航。", GREEN)
    add_card(slide, 4.85, 5.05, 3.2, 0.95, "天气信息", "地点关联 NIWA / MetService weather URL，前端以面板嵌入展示。", BLUE)
    add_card(slide, 8.6, 5.05, 3.2, 0.95, "运行配置", "ALLOWED_ORIGINS、COMMENT_UPLOAD_DIR、数据库连接等通过环境变量控制。", ORANGE)

    # 11. Engineering quality
    slide = new_slide(prs, 11, "工程质量与可维护性", "项目已经按前后端职责拆分，并对关键 API 行为补充后端测试。", "10 QUALITY")
    add_card(slide, 0.85, 1.7, 3.45, 1.05, "测试覆盖", "backend/tests 覆盖 tracks、search、weather、comments、favourites 等接口行为。", GREEN)
    add_card(slide, 4.9, 1.7, 3.45, 1.05, "模块边界", "Repository 只关心 SQL，Service 组合业务搜索，前端组件围绕 UI 场景拆分。", BLUE)
    add_card(slide, 8.95, 1.7, 3.05, 1.05, "数据校验", "Pydantic 限制分页、评分、邮箱长度，并确保评论有文本或图片。", ORANGE)
    add_bullets(
        slide,
        1.1,
        3.45,
        10.8,
        1.65,
        [
            "前端：React 19、TypeScript、Vite、Bootstrap、Leaflet、Recharts；脚本支持 build / lint / preview。",
            "后端：FastAPI、SQLAlchemy、GeoAlchemy2、GeoPandas、Rasterio、Shapely、RapidFuzz、Psycopg。",
            "地理处理：用 CRS 转换、栅格采样和 PostGIS 几何函数保证空间计算可复用。",
            "用户数据：收藏和评论单独 repository 管理，便于后续扩展登录、审核和通知功能。",
        ],
        12,
    )
    add_flow(slide, ["数据可信", "接口清晰", "交互直观", "部署可用"], 1.8, 5.72, 9.7, 0.72, [GREEN, BLUE, ORANGE, DARK_GREEN])

    # 12. Summary
    slide = new_slide(prs, 12, "总结与后续优化", "KiwiTrail 的价值在于把空间数据工程能力转化成用户可直接使用的户外规划体验。", "11 WRAP-UP")
    add_card(slide, 0.9, 1.75, 3.45, 1.25, "项目价值", "集中展示 tracks、huts、campsites、weather stations，降低户外信息检索和行程判断成本。", GREEN)
    add_card(slide, 4.95, 1.75, 3.45, 1.25, "技术亮点", "PostGIS 空间查询、Gazetteer 模糊搜索、DEM 海拔采样、React 地图交互形成完整闭环。", BLUE)
    add_card(slide, 9.0, 1.75, 3.05, 1.25, "可展示成果", "在线 Demo、项目截图、可运行前后端、测试用例和数据处理脚本。", ORANGE)
    add_bullets(
        slide,
        1.15,
        3.75,
        10.3,
        1.45,
        [
            "下一步 1：完善移动端细节与离线/弱网体验，提升户外现场可用性。",
            "下一步 2：扩展路线推荐逻辑，引入距离、累计爬升、天气风险和用户偏好。",
            "下一步 3：增强数据更新自动化，形成定期 ETL、监控和可回滚的数据发布流程。",
        ],
        13,
    )
    thanks = slide.shapes.add_textbox(Inches(0.95), Inches(6.08), Inches(11.2), Inches(0.45))
    p = thanks.text_frame.paragraphs[0]
    p.text = "Thank you · Questions?"
    p.alignment = PP_ALIGN.CENTER
    set_run(p.runs[0], 20, DARK_GREEN, True, FONT_EN)

    prs.save(OUT)


if __name__ == "__main__":
    build_presentation()
    print(f"Created {OUT}")
